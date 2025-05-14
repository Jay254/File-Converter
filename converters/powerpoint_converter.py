import os
from pptx import Presentation
import comtypes.client
import pythoncom

def convert_powerpoint(filepath, target_format):
    """
    Convert PowerPoint files to various formats
    """
    output_path = os.path.splitext(filepath)[0] + f'.{target_format}'
    
    if target_format == 'pdf':
        # PowerPoint to PDF conversion
        powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
        powerpoint.Visible = True
        try:
            deck = powerpoint.Presentations.Open(os.path.abspath(filepath))
            deck.SaveAs(os.path.abspath(output_path), 32)  # 32 is the PDF format
        finally:
            deck.Close()
            powerpoint.Quit()
            
    elif target_format == 'jpg':
        # PowerPoint to JPG conversion (first slide only)
        powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
        powerpoint.Visible = True
        try:
            deck = powerpoint.Presentations.Open(os.path.abspath(filepath))
            deck.Slides[0].Export(os.path.abspath(output_path), 'JPG')
        finally:
            deck.Close()
            powerpoint.Quit()
            
    elif target_format == 'html':
        # PowerPoint to HTML conversion
        prs = Presentation(filepath)
        html_content = ['<html><body>']
        
        for slide in prs.slides:
            html_content.append('<div class="slide">')
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    html_content.append(f'<p>{shape.text}</p>')
            html_content.append('</div>')
            
        html_content.append('</body></html>')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))
            
    else:
        raise ValueError(f"Conversion from PowerPoint to {target_format} is not supported")
        
    return output_path 